import pptx
def pptxwordsReplacer (pptx_path, keyword_dict, output_path):
    prs = pptx.Presentation(pptx_path)
    slovar = keyword_dict
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
        for k, v in slovar.items():
            if k in text_frame.text:
                text_frame.text = text_frame.text.replace(k, v)
    prs.save(output_path)


keyword_dict = {
    "^^FIO^^": "Иванов Иван Иванович",
    "&(Professia)&": "Кондитер"
}

inputfile = "Сертификатик.pptx"
outputfile = "Сертификат.pptx"
pptxwordsReplacer(inputfile, keyword_dict, outputfile)
